"""Generate the bundled example.docx and example.pptx templates.

Run this once and commit the outputs.  The generated files are loaded
by users via :func:`dash_reportbuilder.example_template_path`.

Usage:
    cd dash-reportbuilder
    uv run --active python scripts/build_example_templates.py
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu
from pptx.util import Inches as PPTXInches
from pptx.util import Pt as PPTXPt

# ---------------------------------------------------------------------------
# Brand definition (kept in lockstep with example.typ / example.html)
# ---------------------------------------------------------------------------

BRAND_PRIMARY = RGBColor(0x25, 0x63, 0xEB)  # royal blue
BRAND_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)  # amber
BRAND_MUTED = RGBColor(0x64, 0x74, 0x8B)  # slate
BRAND_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)  # very light gray
BODY_FONT = "Calibri"

# PowerPoint variants of the same colors (separate type).
PPTX_PRIMARY = PPTXColor(0x25, 0x63, 0xEB)
PPTX_ACCENT = PPTXColor(0xF5, 0x9E, 0x0B)
PPTX_MUTED = PPTXColor(0x64, 0x74, 0x8B)

OUT_DIR = Path(__file__).resolve().parent.parent / "src/dash_reportbuilder/templates"


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------


def build_docx() -> Path:
    doc = Document()

    # Body style
    body = doc.styles["Normal"]
    body.font.name = BODY_FONT
    body.font.size = Pt(11)

    # Heading styles — H1 is large + brand color, H2 is medium + brand,
    # H3 is muted gray.
    for lvl, size, color in [
        (1, 22, BRAND_PRIMARY),
        (2, 14, BRAND_PRIMARY),
        (3, 12, BRAND_MUTED),
    ]:
        s = doc.styles[f"Heading {lvl}"]
        s.font.name = BODY_FONT
        s.font.size = Pt(size)
        s.font.bold = True
        s.font.color.rgb = color

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Page header — small slate-gray "Report" line on the right
    header = doc.sections[0].header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.add_run("Report")
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = BRAND_MUTED

    # Page footer — page number on the right
    footer = doc.sections[0].footer
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = fp.add_run("")
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = BRAND_MUTED

    # Caption style — used by DocxBackend.add_caption
    if "Caption" not in [s.name for s in doc.styles]:
        cap = doc.styles.add_style("Caption", WD_STYLE_TYPE.PARAGRAPH)
    else:
        cap = doc.styles["Caption"]
    cap.font.name = BODY_FONT
    cap.font.size = Pt(10)
    cap.font.italic = True
    cap.font.color.rgb = BRAND_MUTED

    # Table base style — header row gets brand fill via existing built-in
    # style.  We just make sure it exists; DocxBackend.add_table sets
    # table.style = "Light Grid Accent 1" itself.

    out = OUT_DIR / "example.docx"
    doc.save(str(out))
    print(f"wrote {out}")
    return out


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------


def _add_brand_bar(slide, prs):
    """Add a thin brand-colored bar at the top of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Emu(0),
        top=Emu(0),
        width=prs.slide_width,
        height=PPTXInches(0.25),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_PRIMARY
    bar.line.fill.background()
    return bar


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = PPTXInches(13.333)
    prs.slide_height = PPTXInches(7.5)

    # Brand all title placeholders across layouts: brand-blue, bold, Calibri
    for layout in prs.slide_layouts:
        for shape in layout.placeholders:
            if not shape.has_text_frame:
                continue
            if shape.placeholder_format.idx == 0:  # title
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = BODY_FONT
                        run.font.color.rgb = PPTX_PRIMARY
                        run.font.bold = True

    # Add a sample title slide so the .pptx isn't blank when opened.
    # Users who pass this template via PptxTemplate(template_path=...) will
    # load this slide; appended slides come after.
    title_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(title_layout)
    _add_brand_bar(slide, prs)

    if slide.shapes.title is not None:
        slide.shapes.title.text = "Report"
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = BODY_FONT
                run.font.color.rgb = PPTX_PRIMARY
                run.font.bold = True
                run.font.size = PPTXPt(40)

    # Subtitle (placeholder idx 1 on Title Slide)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1 and shape.has_text_frame:
            shape.text_frame.text = "Generated with dash-reportbuilder"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = BODY_FONT
                    run.font.color.rgb = PPTX_MUTED
                    run.font.size = PPTXPt(20)

    out = OUT_DIR / "example.pptx"
    prs.save(str(out))
    print(f"wrote {out}")
    return out


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()
