# Copyright (c) Simon Niederberger.
# Distributed under the terms of the MIT License.

"""PowerPoint (.pptx) export via python-pptx."""

from __future__ import annotations

import io

from dash_reportbuilder.export._base import PptxTemplate, decode_data_uri
from dash_reportbuilder.model import ItemType, Report


def export_pptx(report: Report, template: PptxTemplate | None = None) -> bytes:
    """Export *report* to .pptx bytes.

    Each image gets its own slide.  Text items between images are added
    as text boxes on the same slide.

    Requires ``python-pptx`` (install with ``pip install dash-reportbuilder[pptx]``).
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches, Pt
    except ImportError as e:
        raise ImportError(
            "python-pptx is required for PowerPoint export. "
            "Install with: pip install dash-reportbuilder[pptx]"
        ) from e

    t = template or PptxTemplate()

    if t.template_path:
        prs = Presentation(t.template_path)
    else:
        prs = Presentation()

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def _get_blank_layout():
        layouts = prs.slide_layouts
        idx = min(t.image_layout_index, len(layouts) - 1)
        return layouts[idx]

    current_slide = None

    for item in report.items:
        if item.type == ItemType.IMAGE:
            slide = prs.slides.add_slide(_get_blank_layout())
            current_slide = slide

            img_bytes = decode_data_uri(item.content)
            buf = io.BytesIO(img_bytes)

            # Center the image on the slide
            max_w = slide_width - Inches(1)
            max_h = slide_height - Inches(2)

            # Add image and let python-pptx figure out dimensions
            pic = slide.shapes.add_picture(
                buf, Inches(0.5), Inches(0.5), width=max_w
            )
            # Scale down if too tall
            if pic.height > max_h:
                ratio = max_h / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = max_h
            # Center horizontally
            pic.left = (slide_width - pic.width) // 2
            pic.top = Inches(0.5)

        elif item.type == ItemType.HEADING:
            slide = prs.slides.add_slide(_get_blank_layout())
            current_slide = slide
            from pptx.util import Inches, Pt

            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.content
            p.font.size = Pt(28)
            p.font.bold = True

        elif item.type in (ItemType.PARAGRAPH, ItemType.CAPTION):
            if current_slide is None:
                slide = prs.slides.add_slide(_get_blank_layout())
                current_slide = slide
            txBox = current_slide.shapes.add_textbox(
                Inches(0.5),
                slide_height - Inches(1.5),
                slide_width - Inches(1),
                Inches(1),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.content
            p.font.size = Pt(12)
            if item.type == ItemType.CAPTION:
                p.font.italic = True

        elif item.type == ItemType.PAGE_BREAK:
            current_slide = None

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
