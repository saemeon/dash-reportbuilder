# Copyright (c) Simon Niederberger.
# Distributed under the terms of the MIT License.

"""PowerPoint (.pptx) backend via python-pptx."""

from __future__ import annotations

import io
from collections.abc import Callable
from typing import TYPE_CHECKING

from dash_reportbuilder.export._base import PptxTemplate, decode_data_uri

if TYPE_CHECKING:
    from pptx.presentation import Presentation


class PptxBackend:
    """Builds a PowerPoint presentation incrementally.

    Slide model: each image and each heading start a new slide.
    Paragraphs and captions are appended to the current slide.
    A page break clears the "current slide" pointer so the next text
    item starts a new slide.

    Elements that need full python-pptx control can use :meth:`modify`.
    """

    def __init__(self, template: PptxTemplate | None = None) -> None:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError(
                "python-pptx is required for PowerPoint export. "
                "Install with: pip install dash-reportbuilder[pptx]"
            ) from e

        self.template = template or PptxTemplate()
        if self.template.template_path:
            self.prs = Presentation(self.template.template_path)
        else:
            self.prs = Presentation()

        self._current_slide = None

    def _blank_layout(self):
        layouts = self.prs.slide_layouts
        idx = min(self.template.image_layout_index, len(layouts) - 1)
        return layouts[idx]

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self._blank_layout())
        self._current_slide = slide
        return slide

    def _ensure_slide(self):
        if self._current_slide is None:
            return self._new_slide()
        return self._current_slide

    # ------------------------------------------------------------------
    # Generic primitives
    # ------------------------------------------------------------------

    def add_image(
        self,
        data_uri: str,
        *,
        title: str | None = None,
        caption: str | None = None,
        width_mm: float | None = None,
    ) -> None:
        from pptx.util import Inches, Mm

        slide = self._new_slide()
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        assert slide_width is not None
        assert slide_height is not None

        if title:
            from pptx.util import Pt

            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(0.7)
            )
            p = text_box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            top = Inches(1.2)
        else:
            top = Inches(0.5)

        img_bytes = decode_data_uri(data_uri)
        buf = io.BytesIO(img_bytes)

        max_w = Mm(width_mm) if width_mm else slide_width - Inches(1)
        max_h = slide_height - top - Inches(1.5)

        pic = slide.shapes.add_picture(buf, Inches(0.5), top, width=max_w)
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = max_h
        pic.left = (self.prs.slide_width - pic.width) // 2
        pic.top = top

        if caption:
            self.add_caption(caption)

    def add_heading(self, text: str, level: int = 1) -> None:
        from pptx.util import Inches, Pt

        slide = self._new_slide()
        slide_width = self.prs.slide_width
        assert slide_width is not None
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1)
        )
        p = text_box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(max(28 - (level - 1) * 4, 14))
        p.font.bold = True

    def add_paragraph(self, text: str) -> None:
        from pptx.util import Inches, Pt

        slide = self._ensure_slide()
        slide_height = self.prs.slide_height
        slide_width = self.prs.slide_width
        assert slide_height is not None
        assert slide_width is not None
        text_box = slide.shapes.add_textbox(
            Inches(0.5),
            slide_height - Inches(1.5),
            slide_width - Inches(1),
            Inches(1),
        )
        p = text_box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)

    def add_caption(self, text: str) -> None:
        from pptx.util import Inches, Pt

        slide = self._ensure_slide()
        slide_height = self.prs.slide_height
        slide_width = self.prs.slide_width
        assert slide_height is not None
        assert slide_width is not None
        text_box = slide.shapes.add_textbox(
            Inches(0.5),
            slide_height - Inches(1.0),
            slide_width - Inches(1),
            Inches(0.5),
        )
        p = text_box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.italic = True

    def add_table(self, headers: list[str], rows: list[list[str]]) -> None:
        from pptx.util import Inches

        slide = self._new_slide()
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        assert slide_width is not None
        assert slide_height is not None
        n_cols = len(headers)
        n_rows = 1 + len(rows)
        left = Inches(0.5)
        top = Inches(0.5)
        width = slide_width - Inches(1)
        height = slide_height - Inches(1)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                table.cell(r, c).text = str(value)

    def add_page_break(self) -> None:
        self._current_slide = None

    # ------------------------------------------------------------------
    # Native escape hatch
    # ------------------------------------------------------------------

    def modify(self, fn: Callable[[Presentation], None]) -> None:
        """Run *fn* against the underlying python-pptx ``Presentation``."""
        fn(self.prs)

    # ------------------------------------------------------------------
    # Finalize
    # ------------------------------------------------------------------

    def build(self) -> bytes:
        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()
