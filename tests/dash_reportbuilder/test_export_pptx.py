# Copyright (c) Simon Niederberger.
# Distributed under the terms of the MIT License.

import io

import pytest

from dash_reportbuilder.model import ItemType, Report, ReportItem

_TINY_PNG_URI = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADElEQVR4nGP4z8AAAAMBAQDJ/pLvAAAAAElFTkSuQmCC"


@pytest.fixture
def slides_report():
    report = Report(title="Slides")
    report.append(ReportItem(type=ItemType.HEADING, content="Title Slide"))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.CAPTION, content="Chart 1"))
    report.append(ReportItem(type=ItemType.PAGE_BREAK))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    return report


def _load_pptx(data: bytes):
    from pptx import Presentation

    return Presentation(io.BytesIO(data))


def test_export_pptx_produces_bytes(slides_report):
    from dash_reportbuilder.export._pptx import export_pptx

    result = export_pptx(slides_report)
    assert isinstance(result, bytes)
    assert len(result) > 100
    assert result[:2] == b"PK"


def test_export_pptx_empty_report():
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report()
    result = export_pptx(report)
    assert isinstance(result, bytes)


# --- Additional tests ---


def test_heading_creates_new_slide():
    """A heading item creates its own slide with a text box."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="Heading Test")
    report.append(ReportItem(type=ItemType.HEADING, content="My Title"))

    result = export_pptx(report)
    prs = _load_pptx(result)
    assert len(prs.slides) == 1
    # The slide should contain a textbox with the heading text
    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert "My Title" in texts


def test_multiple_images_create_multiple_slides():
    """Each image gets its own slide."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="Multi Image")
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))

    result = export_pptx(report)
    prs = _load_pptx(result)
    assert len(prs.slides) == 3


def test_page_break_resets_slide():
    """Page break sets current_slide to None, so next text item starts a new slide."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="Break Test")
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.CAPTION, content="Caption on image slide"))
    report.append(ReportItem(type=ItemType.PAGE_BREAK))
    report.append(ReportItem(type=ItemType.PARAGRAPH, content="After break"))

    result = export_pptx(report)
    prs = _load_pptx(result)
    # Slide 1: image + caption, slide 2: paragraph (because page_break reset)
    assert len(prs.slides) == 2


def test_caption_on_current_slide():
    """Caption/paragraph text goes on the current slide, not a new one."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="Caption Test")
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.CAPTION, content="My Caption"))

    result = export_pptx(report)
    prs = _load_pptx(result)
    # Only 1 slide: image + caption on the same slide
    assert len(prs.slides) == 1
    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert "My Caption" in texts


def test_paragraph_without_prior_slide_creates_one():
    """A paragraph with no prior slide creates a new slide."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="Orphan Paragraph")
    report.append(ReportItem(type=ItemType.PARAGRAPH, content="Standalone text"))

    result = export_pptx(report)
    prs = _load_pptx(result)
    assert len(prs.slides) == 1


def test_heading_then_image_creates_two_slides():
    """Heading and image each get their own slide."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report(title="H+I")
    report.append(ReportItem(type=ItemType.HEADING, content="Heading"))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))

    result = export_pptx(report)
    prs = _load_pptx(result)
    assert len(prs.slides) == 2


def test_caption_italic():
    """Caption text is italic."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report()
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))
    report.append(ReportItem(type=ItemType.CAPTION, content="Italic caption"))

    result = export_pptx(report)
    prs = _load_pptx(result)
    # Find the caption text box
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text == "Italic caption":
            assert shape.text_frame.paragraphs[0].font.italic is True
            break
    else:
        pytest.fail("Caption text box not found")


def test_slides_report_slide_count(slides_report):
    """slides_report: heading(1 slide) + image(1 slide) + caption(on image slide) + page_break + image(1 slide) = 3 slides."""
    from dash_reportbuilder.export._pptx import export_pptx

    result = export_pptx(slides_report)
    prs = _load_pptx(result)
    assert len(prs.slides) == 3


def test_multiple_page_breaks():
    """Multiple consecutive page breaks should not crash."""
    from dash_reportbuilder.export._pptx import export_pptx

    report = Report()
    report.append(ReportItem(type=ItemType.PAGE_BREAK))
    report.append(ReportItem(type=ItemType.PAGE_BREAK))
    report.append(ReportItem(type=ItemType.IMAGE, content=_TINY_PNG_URI))

    result = export_pptx(report)
    prs = _load_pptx(result)
    # Only the image creates a slide; page breaks alone create no slides
    assert len(prs.slides) == 1
