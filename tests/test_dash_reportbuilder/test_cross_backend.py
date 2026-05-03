# Copyright (c) Simon Niederberger.
# Distributed under the terms of the MIT License.

"""Tests that span every backend.

Catches "I added a new method to one backend and forgot the others" and
locks down the FileStore-save → reload → export-via-each-backend pipeline
that real users hit.
"""

from __future__ import annotations

import io
import tempfile
import zipfile
from collections.abc import Callable

import pytest

from dash_reportbuilder import (
    CaptionElement,
    DocxBackend,
    HeadingElement,
    HtmlBackend,
    ImageElement,
    ImageZipBackend,
    PageBreakElement,
    ParagraphElement,
    PptxBackend,
    Report,
    TypstBackend,
)
from dash_reportbuilder.protocols import ReportBackend
from dash_reportbuilder.store import FileStore

TINY_PNG_URI = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADElEQVR4nGP4z8AAAAMBAQDJ/pLvAAAAAElFTkSuQmCC"


def _build_report() -> Report:
    """Build a report with one of every element type, including a titled image."""
    report = Report(title="Cross-Backend Report")
    report.append(HeadingElement(text="Section One", level=2))
    report.append(ParagraphElement(text="Some prose body."))
    report.append(ImageElement(data_uri=TINY_PNG_URI, title="MyChart"))
    report.append(CaptionElement(text="Figure 1"))
    report.append(PageBreakElement())
    report.append(HeadingElement(text="Section Two", level=2))
    report.append(ParagraphElement(text="More prose."))
    return report


# ---------------------------------------------------------------------------
# Backend factories
# ---------------------------------------------------------------------------


def _docx(report: Report) -> bytes:
    return report.export(DocxBackend(title=report.title))


def _pptx(report: Report) -> bytes:
    return report.export(PptxBackend())


def _typst_source(report: Report) -> bytes:
    backend = TypstBackend(title=report.title)
    for it in report.items:
        it.render_into(backend)
    return backend.build_source().encode("utf-8")


def _html(report: Report) -> bytes:
    return report.export(HtmlBackend(title=report.title))


def _images(report: Report) -> bytes:
    return report.export(ImageZipBackend())


_FACTORIES: dict[str, Callable[[Report], bytes]] = {
    "docx": _docx,
    "pptx": _pptx,
    "typst": _typst_source,
    "html": _html,
    "images": _images,
}


# ---------------------------------------------------------------------------
# (2) Cross-backend consistency
# ---------------------------------------------------------------------------


class TestEveryBackendProducesNonEmptyOutput:
    """Every backend produces non-empty bytes.

    The same Report runs through every backend without crashing and
    produces non-empty bytes.
    """

    @pytest.mark.parametrize("name", list(_FACTORIES))
    def test_non_empty(self, name):
        report = _build_report()
        data = _FACTORIES[name](report)
        assert isinstance(data, bytes)
        assert len(data) > 0


class TestTextContentReachesEachBackend:
    """Text content reaches every text backend.

    User-facing strings (heading, paragraph) appear in document
    backends.  Image-only and binary backends are checked separately.
    """

    @pytest.mark.parametrize(
        "name,decode",
        [
            ("typst", lambda b: b.decode("utf-8")),
            ("html", lambda b: b.decode("utf-8")),
        ],
    )
    def test_text_in_text_backends(self, name, decode):
        report = _build_report()
        text = decode(_FACTORIES[name](report))
        assert "Section One" in text
        assert "Section Two" in text
        assert "Some prose body." in text
        assert "Figure 1" in text

    def test_text_in_docx(self):
        from docx import Document

        report = _build_report()
        data = _FACTORIES["docx"](report)
        doc = Document(io.BytesIO(data))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "Section One" in all_text
        assert "Some prose body." in all_text
        assert "Figure 1" in all_text

    def test_text_in_pptx(self):
        from pptx import Presentation

        report = _build_report()
        data = _FACTORIES["pptx"](report)
        prs = Presentation(io.BytesIO(data))
        texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        flat = "\n".join(texts)
        assert "Section One" in flat
        assert "Some prose body." in flat


class TestImageReachesEachBackend:
    """Every backend that supports images includes the image somehow."""

    def test_html_embeds_data_uri(self):
        text = _FACTORIES["html"](_build_report()).decode("utf-8")
        assert TINY_PNG_URI in text

    def test_typst_references_image_filename(self):
        text = _FACTORIES["typst"](_build_report()).decode("utf-8")
        assert 'image("img_1.png"' in text

    def test_docx_has_one_inline_shape(self):
        from docx import Document

        data = _FACTORIES["docx"](_build_report())
        doc = Document(io.BytesIO(data))
        assert len(doc.inline_shapes) == 1

    def test_pptx_has_a_picture_shape(self):
        from pptx import Presentation

        data = _FACTORIES["pptx"](_build_report())
        prs = Presentation(io.BytesIO(data))
        n_pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        assert n_pics >= 1

    def test_image_zip_contains_titled_file(self):
        data = _FACTORIES["images"](_build_report())
        names = zipfile.ZipFile(io.BytesIO(data)).namelist()
        assert "MyChart.png" in names


class TestBackendsImplementProtocol:
    """Every backend implements the ReportBackend Protocol.

    Every backend exposes the full primitive surface declared by the
    ReportBackend Protocol.
    """

    @pytest.mark.parametrize(
        "backend",
        [
            DocxBackend(),
            PptxBackend(),
            TypstBackend(),
            HtmlBackend(),
            ImageZipBackend(),
        ],
    )
    def test_has_all_primitives(self, backend):
        for method in (
            "add_image",
            "add_heading",
            "add_paragraph",
            "add_caption",
            "add_table",
            "add_page_break",
            "build",
        ):
            assert callable(getattr(backend, method)), (
                f"{type(backend).__name__} is missing {method}"
            )

    @pytest.mark.parametrize(
        "backend",
        [
            DocxBackend(),
            PptxBackend(),
            TypstBackend(),
            HtmlBackend(),
            ImageZipBackend(),
        ],
    )
    def test_satisfies_protocol(self, backend):
        # ReportBackend is a runtime-checkable Protocol.
        assert isinstance(backend, ReportBackend)


# ---------------------------------------------------------------------------
# (3) FileStore + export pipeline
# ---------------------------------------------------------------------------


class TestFileStoreReloadAndExport:
    """Reload then export matches in-memory export.

    Save a report, reload from a fresh FileStore, export — every
    backend produces output equivalent to the in-memory path.
    """

    @pytest.mark.parametrize("name", list(_FACTORIES))
    def test_reloaded_export_matches_in_memory(self, name):
        report = _build_report()
        in_memory = _FACTORIES[name](report)

        with tempfile.TemporaryDirectory() as tmpdir:
            store = FileStore(tmpdir)
            store.put("session", report)
            loaded = FileStore(tmpdir).get("session")

        from_disk = _FACTORIES[name](loaded)

        # Binary backends (.docx, .pptx, .zip) embed timestamps in their
        # zip metadata, so byte-for-byte equality won't hold.  We check
        # the comparable structural invariants instead.
        if name == "docx":
            from docx import Document

            d_mem = Document(io.BytesIO(in_memory))
            d_disk = Document(io.BytesIO(from_disk))
            assert [p.text for p in d_mem.paragraphs] == [
                p.text for p in d_disk.paragraphs
            ]
            assert len(d_mem.inline_shapes) == len(d_disk.inline_shapes)
        elif name == "pptx":
            from pptx import Presentation

            p_mem = Presentation(io.BytesIO(in_memory))
            p_disk = Presentation(io.BytesIO(from_disk))
            assert len(p_mem.slides) == len(p_disk.slides)
        elif name == "images":
            assert sorted(zipfile.ZipFile(io.BytesIO(in_memory)).namelist()) == sorted(
                zipfile.ZipFile(io.BytesIO(from_disk)).namelist()
            )
        else:
            # Text backends — bytes are deterministic.
            assert in_memory == from_disk

    def test_reorder_persists_through_reload(self):
        """Reorder before save → reload → export reflects new order."""
        report = _build_report()
        item_ids = [it.id for it in report.items]
        reversed_ids = list(reversed(item_ids))
        report.reorder(reversed_ids)

        with tempfile.TemporaryDirectory() as tmpdir:
            store = FileStore(tmpdir)
            store.put("session", report)
            loaded = FileStore(tmpdir).get("session")

        assert [it.id for it in loaded.items] == reversed_ids

        # And the reordered text appears in HTML in the right order.
        html = _FACTORIES["html"](loaded).decode("utf-8")
        idx_one = html.index("Section One")
        idx_two = html.index("Section Two")
        # After reversal, Section Two should appear before Section One.
        assert idx_two < idx_one
