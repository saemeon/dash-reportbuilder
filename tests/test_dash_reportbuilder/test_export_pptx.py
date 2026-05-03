# Copyright (c) Simon Niederberger.
# Distributed under the terms of the MIT License.

import io

import pytest

from dash_reportbuilder.backends import PptxBackend
from dash_reportbuilder.elements import (
    CaptionElement,
    HeadingElement,
    ImageElement,
    PageBreakElement,
    ParagraphElement,
)
from dash_reportbuilder.model import Report

_TINY_PNG_URI = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADElEQVR4nGP4z8AAAAMBAQDJ/pLvAAAAAElFTkSuQmCC"


@pytest.fixture
def slides_report():
    report = Report(title="Slides")
    report.append(HeadingElement(text="Title Slide", level=2))
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(CaptionElement(text="Chart 1"))
    report.append(PageBreakElement())
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    return report


def _load_pptx(data: bytes):
    from pptx import Presentation

    return Presentation(io.BytesIO(data))


def test_export_pptx_produces_bytes(slides_report):

    result = slides_report.export(PptxBackend())
    assert isinstance(result, bytes)
    assert len(result) > 100
    assert result[:2] == b"PK"


def test_export_pptx_empty_report():

    report = Report()
    result = report.export(PptxBackend())
    assert isinstance(result, bytes)


# --- Additional tests ---


def test_heading_creates_new_slide():
    """A heading item creates its own slide with a text box."""
    report = Report(title="Heading Test")
    report.append(HeadingElement(text="My Title", level=2))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 1
    texts = [
        shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
    ]
    assert "My Title" in texts


def test_multiple_images_create_multiple_slides():
    """Each image gets its own slide."""
    report = Report(title="Multi Image")
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(ImageElement(data_uri=_TINY_PNG_URI))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 3


def test_page_break_resets_slide():
    """Page break sets current_slide to None, so next text item starts a new slide."""
    report = Report(title="Break Test")
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(CaptionElement(text="Caption on image slide"))
    report.append(PageBreakElement())
    report.append(ParagraphElement(text="After break"))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 2


def test_caption_on_current_slide():
    """Caption/paragraph text goes on the current slide, not a new one."""
    report = Report(title="Caption Test")
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(CaptionElement(text="My Caption"))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 1
    texts = [
        shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
    ]
    assert "My Caption" in texts


def test_paragraph_without_prior_slide_creates_one():
    """A paragraph with no prior slide creates a new slide."""
    report = Report(title="Orphan Paragraph")
    report.append(ParagraphElement(text="Standalone text"))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 1


def test_heading_then_image_creates_two_slides():
    """Heading and image each get their own slide."""
    report = Report(title="H+I")
    report.append(HeadingElement(text="Heading", level=2))
    report.append(ImageElement(data_uri=_TINY_PNG_URI))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 2


def test_caption_italic():
    """Caption text is italic."""
    report = Report()
    report.append(ImageElement(data_uri=_TINY_PNG_URI))
    report.append(CaptionElement(text="Italic caption"))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text == "Italic caption":
            assert shape.text_frame.paragraphs[0].font.italic is True
            break
    else:
        pytest.fail("Caption text box not found")


def test_slides_report_slide_count(slides_report):
    """slides_report: heading(1 slide) + image(1 slide) + caption(on image slide) + page_break + image(1 slide) = 3 slides."""
    result = slides_report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 3


def test_multiple_page_breaks():
    """Multiple consecutive page breaks should not crash."""
    report = Report()
    report.append(PageBreakElement())
    report.append(PageBreakElement())
    report.append(ImageElement(data_uri=_TINY_PNG_URI))

    result = report.export(PptxBackend())
    prs = _load_pptx(result)
    assert len(prs.slides) == 1
